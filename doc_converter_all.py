#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Office ve Diğer Belgeleri PDF'e Dönüştürme Betiği
Desteklenen formatlar: DOCX, XLSX, PPTX, CSV, TXT, HTML, RTF ve PDF
"""

import sys
import os
import base64
import json
import tempfile
import io
from fpdf import FPDF
import pandas as pd
import traceback

# Gerekli kütüphaneleri yükle
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    import pypandoc
    PANDOC_AVAILABLE = True
except ImportError:
    PANDOC_AVAILABLE = False

try:
    from PyPDF2 import PdfReader, PdfWriter
    PYPDF2_AVAILABLE = True
except ImportError:
    try:
        from PyPDF2 import PdfFileReader as PdfReader, PdfFileWriter as PdfWriter
        PYPDF2_AVAILABLE = True
    except ImportError:
        PYPDF2_AVAILABLE = False


# PDF oluşturucu sınıf
class PDFConverter:
    def __init__(self):
        self.pdf = FPDF()
        self.pdf.add_page()
        self.pdf.set_font("Arial", size=12)
        self.pdf.set_auto_page_break(auto=True, margin=15)
    
    def add_title(self, title, size=16):
        self.pdf.set_font("Arial", 'B', size=size)
        self.pdf.cell(0, 10, title, ln=True)
        self.pdf.ln(5)
        self.pdf.set_font("Arial", size=12)
    
    def add_text(self, text):
        if text and text.strip():
            # UTF-8 kodlamasıyla metni ekle
            try:
                safe_text = text.encode('latin-1', 'replace').decode('latin-1')
                self.pdf.multi_cell(0, 10, safe_text)
            except Exception as e:
                print(f"Metin dönüştürme hatası: {e}")
                self.pdf.multi_cell(0, 10, "< Dönüştürme hatası >")
    
    def add_table(self, data, headers=None):
        if headers:
            self.pdf.set_font("Arial", 'B', size=12)
            for header in headers:
                self.pdf.cell(40, 10, str(header)[:15], border=1)
            self.pdf.ln()
            self.pdf.set_font("Arial", size=10)
        
        for row in data:
            for cell in row:
                cell_str = str(cell)
                if len(cell_str) > 15:
                    cell_str = cell_str[:12] + "..."
                self.pdf.cell(40, 10, cell_str, border=1)
            self.pdf.ln()
    
    def get_buffer(self):
        return self.pdf.output(dest='S').encode('latin-1')


# TXT dosyasını PDF'e dönüştür
def txt_to_pdf(input_data):
    try:
        # Metin içeriğini çıkar
        text_content = input_data.decode('utf-8', errors='replace')
        
        # PDF oluştur
        converter = PDFConverter()
        converter.add_title("Metin Dosyası", 16)
        
        # Metni satır satır ekle
        for line in text_content.split('\n'):
            converter.add_text(line)
        
        # PDF'i belleğe aktar
        return converter.get_buffer()
    except Exception as e:
        print(f"TXT Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# DOCX dosyasını PDF'e dönüştür
def docx_to_pdf(input_data):
    if not DOCX_AVAILABLE:
        raise ImportError("python-docx kütüphanesi yüklü değil")
    
    # Geçici dosya oluştur
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_docx:
        temp_docx.write(input_data)
        temp_docx_path = temp_docx.name
    
    try:
        # DOCX dosyasını oku
        document = Document(temp_docx_path)
        
        # PDF oluştur
        converter = PDFConverter()
        
        # Başlık ve paragrafları ekle
        for paragraph in document.paragraphs:
            if paragraph.style.name.startswith('Heading'):
                converter.pdf.set_font("Arial", 'B', size=14)
            else:
                converter.pdf.set_font("Arial", size=12)
            
            converter.add_text(paragraph.text)
        
        # Tabloları ekle (basitleştirilmiş)
        for table in document.tables:
            converter.pdf.ln(5)
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cells.append(cell.text)
                
                converter.pdf.set_font("Arial", size=10)
                for cell_text in cells:
                    safe_text = str(cell_text).encode('latin-1', 'replace').decode('latin-1')
                    converter.pdf.cell(40, 10, safe_text[:15], border=1)
                converter.pdf.ln()
        
        # PDF'i belleğe aktar ve geçici dosyayı temizle
        result = converter.get_buffer()
        os.unlink(temp_docx_path)
        return result
    
    except Exception as e:
        # Hata durumunda temizlik
        if os.path.exists(temp_docx_path):
            os.unlink(temp_docx_path)
        print(f"DOCX Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# XLSX dosyasını PDF'e dönüştür
def xlsx_to_pdf(input_data):
    if not EXCEL_AVAILABLE:
        raise ImportError("openpyxl kütüphanesi yüklü değil")
    
    # Geçici dosya oluştur
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as temp_xlsx:
        temp_xlsx.write(input_data)
        temp_xlsx_path = temp_xlsx.name
    
    try:
        # Excel dosyasını oku
        workbook = openpyxl.load_workbook(temp_xlsx_path)
        
        # PDF oluştur
        converter = PDFConverter()
        
        # Her çalışma sayfası için
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            converter.add_title(f"Çalışma Sayfası: {sheet_name}", 14)
            
            # Verileri tablo olarak ekle
            for row in sheet.iter_rows(values_only=True):
                row_items = []
                for cell in row:
                    row_items.append(str(cell) if cell is not None else "")
                
                converter.pdf.set_font("Arial", size=10)
                for cell_text in row_items:
                    safe_text = cell_text.encode('latin-1', 'replace').decode('latin-1')
                    converter.pdf.cell(40, 10, safe_text[:15], border=1)
                converter.pdf.ln()
            
            # Sayfa sonuna boşluk ekle
            converter.pdf.ln(10)
            
            # Çalışma sayfası bittiyse ve başka sayfa varsa yeni sayfa ekle
            if sheet_name != workbook.sheetnames[-1]:
                converter.pdf.add_page()
        
        # PDF'i belleğe aktar ve geçici dosyayı temizle
        result = converter.get_buffer()
        os.unlink(temp_xlsx_path)
        return result
    
    except Exception as e:
        # Hata durumunda temizlik
        if os.path.exists(temp_xlsx_path):
            os.unlink(temp_xlsx_path)
        print(f"XLSX Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# PPTX dosyasını PDF'e dönüştür
def pptx_to_pdf(input_data):
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx kütüphanesi yüklü değil")
    
    # Geçici dosya oluştur
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx:
        temp_pptx.write(input_data)
        temp_pptx_path = temp_pptx.name
    
    try:
        # PowerPoint dosyasını oku
        presentation = Presentation(temp_pptx_path)
        
        # PDF oluştur
        converter = PDFConverter()
        
        # Her slayt için
        for i, slide in enumerate(presentation.slides):
            # Her slayt için yeni sayfa ekle (ilk sayfa hariç)
            if i > 0:
                converter.pdf.add_page()
            
            converter.add_title(f"Slayt {i+1}", 16)
            
            # Slayt içeriğini ekle
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if shape.text.strip():
                        converter.add_text(shape.text)
        
        # PDF'i belleğe aktar ve geçici dosyayı temizle
        result = converter.get_buffer()
        os.unlink(temp_pptx_path)
        return result
    
    except Exception as e:
        # Hata durumunda temizlik
        if os.path.exists(temp_pptx_path):
            os.unlink(temp_pptx_path)
        print(f"PPTX Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# HTML dosyasını PDF'e dönüştür
def html_to_pdf(input_data):
    if not BS4_AVAILABLE:
        raise ImportError("beautifulsoup4 kütüphanesi yüklü değil")
    
    try:
        # HTML içeriğini çıkar
        html_content = input_data.decode('utf-8', errors='replace')
        
        # BeautifulSoup ile parse et
        soup = BeautifulSoup(html_content, "html.parser")
        
        # PDF oluştur
        converter = PDFConverter()
        
        # Başlık
        title = soup.title.string if soup.title else "HTML Belgesi"
        converter.add_title(title, 16)
        
        # Metin içeriğini ekle
        # Basit yaklaşım: tüm metni ekle
        text_content = soup.get_text(separator='\n', strip=True)
        
        # Metni satır satır ekle
        for line in text_content.split('\n'):
            if line.strip():
                converter.add_text(line)
        
        # PDF'i belleğe aktar
        return converter.get_buffer()
    
    except Exception as e:
        print(f"HTML Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# CSV dosyasını PDF'e dönüştür
def csv_to_pdf(input_data):
    try:
        # CSV içeriğini çıkar
        csv_content = input_data.decode('utf-8', errors='replace')
        csv_file = io.StringIO(csv_content)
        
        # CSV'yi pandas ile oku
        df = pd.read_csv(csv_file)
        
        # PDF oluştur
        converter = PDFConverter()
        converter.add_title("CSV Verileri", 16)
        
        # Başlıkları ekle
        headers = df.columns.tolist()
        converter.pdf.set_font("Arial", 'B', size=12)
        
        # Tablonun genişlik kontrolü
        col_width = min(40, 180 / len(headers)) if headers else 40
        
        # Başlıkları ekle
        for header in headers:
            safe_header = str(header).encode('latin-1', 'replace').decode('latin-1')
            if len(safe_header) > 15:
                safe_header = safe_header[:12] + "..."
            converter.pdf.cell(col_width, 10, safe_header, border=1)
        converter.pdf.ln()
        
        # Verileri ekle
        converter.pdf.set_font("Arial", size=10)
        for _, row in df.iterrows():
            for item in row:
                item_str = str(item)
                safe_text = item_str.encode('latin-1', 'replace').decode('latin-1')
                if len(safe_text) > 15:
                    safe_text = safe_text[:12] + "..."
                converter.pdf.cell(col_width, 10, safe_text, border=1)
            converter.pdf.ln()
        
        # PDF'i belleğe aktar
        return converter.get_buffer()
    
    except Exception as e:
        print(f"CSV Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# RTF dosyasını PDF'e dönüştür
def rtf_to_pdf(input_data):
    if not PANDOC_AVAILABLE:
        raise ImportError("pypandoc kütüphanesi yüklü değil")
    
    # Geçici dosya oluştur
    with tempfile.NamedTemporaryFile(suffix='.rtf', delete=False) as temp_rtf:
        temp_rtf.write(input_data)
        temp_rtf_path = temp_rtf.name
    
    try:
        # RTF dosyasını oku ve dönüştür
        text = pypandoc.convert_file(temp_rtf_path, 'plain')
        
        # PDF oluştur
        converter = PDFConverter()
        converter.add_title("RTF Belgesi", 16)
        
        # Metni satır satır ekle
        for line in text.split('\n'):
            converter.add_text(line)
        
        # PDF'i belleğe aktar ve geçici dosyayı temizle
        result = converter.get_buffer()
        os.unlink(temp_rtf_path)
        return result
    
    except Exception as e:
        # Hata durumunda temizlik
        if os.path.exists(temp_rtf_path):
            os.unlink(temp_rtf_path)
        print(f"RTF Dönüştürme Hatası: {e}")
        traceback.print_exc()
        raise e


# PDF dosyasını kopyala
def pdf_to_pdf(input_data):
    if not PYPDF2_AVAILABLE:
        raise ImportError("PyPDF2 kütüphanesi yüklü değil")
    
    # PDF zaten PDF ise, doğrudan döndür
    return input_data


# Ana dönüştürme fonksiyonu
def convert_to_pdf(input_data, mime_type, file_name):
    """
    Belgeyi türüne göre PDF'e dönüştürür
    
    Args:
        input_data (bytes): Dosya içeriği
        mime_type (str): MIME türü
        file_name (str): Dosya adı (uzantıyı tespit için kullanılır)
    
    Returns:
        bytes: PDF içeriği
    """
    try:
        print(f"Dönüştürülüyor: {file_name}, MIME: {mime_type}")
        
        # Dosya türüne göre uygun dönüştürücüyü çağır
        lower_name = file_name.lower()
        
        # DOCX dönüştürme
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or lower_name.endswith('.docx'):
            return docx_to_pdf(input_data)
        
        # XLSX dönüştürme
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or lower_name.endswith('.xlsx'):
            return xlsx_to_pdf(input_data)
        
        # PPTX dönüştürme
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or lower_name.endswith('.pptx'):
            return pptx_to_pdf(input_data)
        
        # CSV dönüştürme
        elif mime_type == "text/csv" or lower_name.endswith('.csv'):
            return csv_to_pdf(input_data)
        
        # TXT dönüştürme
        elif mime_type == "text/plain" or lower_name.endswith('.txt'):
            return txt_to_pdf(input_data)
        
        # HTML dönüştürme
        elif mime_type == "text/html" or lower_name.endswith('.html') or lower_name.endswith('.htm'):
            return html_to_pdf(input_data)
        
        # RTF dönüştürme
        elif mime_type == "application/rtf" or lower_name.endswith('.rtf'):
            return rtf_to_pdf(input_data)
        
        # PDF'i olduğu gibi döndür
        elif mime_type == "application/pdf" or lower_name.endswith('.pdf'):
            return pdf_to_pdf(input_data)
        
        else:
            raise ValueError(f"Desteklenmeyen dosya türü: {mime_type}")
    
    except Exception as e:
        print(f"Dönüştürme hatası: {str(e)}")
        traceback.print_exc()
        raise e


def main():
    """
    Ana fonksiyon - Node.js'den çağrılacak
    
    Beklenen format:
    1. Argüman: Base64 formatında dosya içeriği
    2. Argüman: MIME türü
    3. Argüman: Dosya adı
    
    Çıktı:
    Base64 formatında PDF içeriği (stdout'a yazılır)
    """
    try:
        if len(sys.argv) < 4:
            raise ValueError("Eksik argümanlar. Beklenen format: <base64_file_content> <mime_type> <file_name>")
        
        # Argümanları al
        base64_content = sys.argv[1]
        mime_type = sys.argv[2]
        file_name = sys.argv[3]
        
        print(f"İşleniyor: {file_name}, MIME: {mime_type}")
        
        # Base64'ü decode et
        file_content = base64.b64decode(base64_content)
        
        # Dönüştürme işlemi
        pdf_bytes = convert_to_pdf(file_content, mime_type, file_name)
        
        # Sonucu base64 olarak encode et ve stdout'a yaz
        pdf_base64 = base64.b64encode(pdf_bytes).decode('utf-8')
        
        # JSON formatında sonuç döndür
        result = {
            "success": True,
            "pdf_base64": pdf_base64,
            "original_size": len(file_content),
            "pdf_size": len(pdf_bytes)
        }
        print(json.dumps(result))
        
    except Exception as e:
        # Hata durumunda hata mesajı döndür
        error_result = {
            "success": False,
            "error": str(e)
        }
        print(json.dumps(error_result))
        sys.exit(1)


if __name__ == "__main__":
    main()